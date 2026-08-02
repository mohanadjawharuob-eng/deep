"""Build the pitch deck as a real PowerPoint file.

    cd backend && .venv/bin/pip install python-pptx
    cd .. && backend/.venv/bin/python docs/pitch/build_deck.py

Written rather than drawn, for one reason: a deck that is a *file* somebody
made once goes stale the week after, and nobody can fix a typo without the
original. This regenerates in a second, lives in the repository next to the
thing it describes, and can be edited by changing a sentence in a list.

``python-pptx`` is deliberately not a dependency of the platform. It is needed
to build this file and by nothing that runs in production, so it is installed
when the deck is rebuilt and not before.

The output is an ordinary .pptx. Open it in PowerPoint, Keynote, Google Slides
or LibreOffice and change anything — it is not locked, and it is not a picture
of a deck.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# The palette, from the platform's own design
# --------------------------------------------------------------------------
CLAY = RGBColor(0x8B, 0x3A, 0x1F)  # fired clay — the accent, spent sparingly
BONE = RGBColor(0xF2, 0xED, 0xE3)  # the page
INK = RGBColor(0x3B, 0x2F, 0x24)  # text
MUTED = RGBColor(0x7A, 0x6B, 0x5C)  # secondary text
BRONZE = RGBColor(0x2F, 0x5D, 0x50)  # oxidised bronze — the second voice
RULE = RGBColor(0xD9, 0xD0, 0xBF)

SERIF = "Georgia"
SANS = "Calibri"
MONO = "Consolas"

WIDTH = Inches(13.333)
HEIGHT = Inches(7.5)
MARGIN = Inches(0.9)


def _background(slide, colour: RGBColor = BONE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _text(
    slide,
    text: str,
    *,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    size: int = 18,
    bold: bool = False,
    colour: RGBColor = INK,
    font: str = SANS,
    align=PP_ALIGN.LEFT,
    spacing: float = 1.0,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0

    for index, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = spacing
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = font
    return box


def _rule(slide, *, top: Emu, left: Emu = MARGIN, width: Emu | None = None, colour=CLAY):
    from pptx.enum.shapes import MSO_SHAPE

    width = width or Inches(1.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _card(slide, *, left, top, width, height, fill=RGBColor(0xFC, 0xFA, 0xF6)):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.04
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RULE
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


def blank(deck: Presentation, *, colour: RGBColor = BONE):
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    _background(slide, colour)
    return slide


# --------------------------------------------------------------------------
# Slide kinds
# --------------------------------------------------------------------------
def title_slide(deck: Presentation, title: str, subtitle: str, footer: str) -> None:
    slide = blank(deck, colour=INK)

    _rule(slide, top=Inches(2.5), left=MARGIN, width=Inches(2.2))
    _text(
        slide,
        title,
        left=MARGIN,
        top=Inches(2.9),
        width=Inches(10),
        height=Inches(1.6),
        size=66,
        bold=True,
        colour=BONE,
        font=SERIF,
    )
    _text(
        slide,
        subtitle,
        left=MARGIN,
        top=Inches(4.4),
        width=Inches(9),
        height=Inches(1.2),
        size=22,
        colour=RGBColor(0xC9, 0xBE, 0xAC),
        spacing=1.3,
    )
    _text(
        slide,
        footer,
        left=MARGIN,
        top=Inches(6.4),
        width=Inches(10),
        height=Inches(0.5),
        size=13,
        colour=RGBColor(0x8E, 0x82, 0x74),
    )


def statement(deck: Presentation, kicker: str, line: str, note: str = "") -> None:
    """One sentence, room around it. For the turns in the argument."""
    slide = blank(deck)

    _text(
        slide,
        kicker.upper(),
        left=MARGIN,
        top=Inches(1.5),
        width=Inches(8),
        height=Inches(0.4),
        size=12,
        bold=True,
        colour=CLAY,
    )
    _rule(slide, top=Inches(2.0), width=Inches(1.2))
    _text(
        slide,
        line,
        left=MARGIN,
        top=Inches(2.5),
        width=Inches(11),
        height=Inches(2.6),
        size=40,
        bold=True,
        colour=INK,
        font=SERIF,
        spacing=1.15,
    )
    if note:
        _text(
            slide,
            note,
            left=MARGIN,
            top=Inches(5.3),
            width=Inches(10),
            height=Inches(1.4),
            size=17,
            colour=MUTED,
            spacing=1.35,
        )


def bullets(deck: Presentation, heading: str, items: list[tuple[str, str]], kicker: str = "") -> None:
    slide = blank(deck)

    top = Inches(0.9)
    if kicker:
        _text(
            slide,
            kicker.upper(),
            left=MARGIN,
            top=top,
            width=Inches(8),
            height=Inches(0.35),
            size=12,
            bold=True,
            colour=CLAY,
        )
        top = Inches(1.35)

    _text(
        slide,
        heading,
        left=MARGIN,
        top=top,
        width=Inches(11),
        height=Inches(0.9),
        size=34,
        bold=True,
        colour=INK,
        font=SERIF,
    )
    _rule(slide, top=top + Inches(0.95), width=Inches(1.2))

    y = top + Inches(1.45)
    for label, body in items:
        _text(
            slide,
            label,
            left=MARGIN,
            top=y,
            width=Inches(11.4),
            height=Inches(0.35),
            size=19,
            bold=True,
            colour=BRONZE,
        )
        _text(
            slide,
            body,
            left=MARGIN,
            top=y + Inches(0.36),
            width=Inches(11.4),
            height=Inches(0.7),
            size=15,
            colour=MUTED,
            spacing=1.3,
        )
        y += Inches(1.12)


def grid(deck: Presentation, heading: str, cells: list[tuple[str, str]], kicker: str = "") -> None:
    """Three across, two down. For the modules."""
    slide = blank(deck)

    if kicker:
        _text(
            slide,
            kicker.upper(),
            left=MARGIN,
            top=Inches(0.85),
            width=Inches(8),
            height=Inches(0.35),
            size=12,
            bold=True,
            colour=CLAY,
        )
    _text(
        slide,
        heading,
        left=MARGIN,
        top=Inches(1.25),
        width=Inches(11),
        height=Inches(0.8),
        size=34,
        bold=True,
        colour=INK,
        font=SERIF,
    )
    _rule(slide, top=Inches(2.15), width=Inches(1.2))

    card_w = Inches(3.5)
    card_h = Inches(1.85)
    gap = Inches(0.32)
    for index, (label, body) in enumerate(cells):
        column = index % 3
        row = index // 3
        left = MARGIN + column * (card_w + gap)
        top = Inches(2.6) + row * (card_h + gap)

        _card(slide, left=left, top=top, width=card_w, height=card_h)
        _text(
            slide,
            label,
            left=left + Inches(0.28),
            top=top + Inches(0.26),
            width=card_w - Inches(0.56),
            height=Inches(0.4),
            size=17,
            bold=True,
            colour=INK,
        )
        _text(
            slide,
            body,
            left=left + Inches(0.28),
            top=top + Inches(0.72),
            width=card_w - Inches(0.56),
            height=Inches(0.95),
            size=12,
            colour=MUTED,
            spacing=1.25,
        )


def compare(deck: Presentation, heading: str, before: list[str], after: list[str]) -> None:
    slide = blank(deck)

    _text(
        slide,
        heading,
        left=MARGIN,
        top=Inches(1.0),
        width=Inches(11),
        height=Inches(0.8),
        size=34,
        bold=True,
        colour=INK,
        font=SERIF,
    )
    _rule(slide, top=Inches(1.9), width=Inches(1.2))

    column_w = Inches(5.3)
    for index, (title, lines, accent) in enumerate(
        [("Today", before, MUTED), ("With Stratum", after, BRONZE)]
    ):
        left = MARGIN + index * (column_w + Inches(0.7))
        _card(slide, left=left, top=Inches(2.4), width=column_w, height=Inches(3.9))
        _text(
            slide,
            title.upper(),
            left=left + Inches(0.35),
            top=Inches(2.7),
            width=column_w - Inches(0.7),
            height=Inches(0.35),
            size=13,
            bold=True,
            colour=accent,
        )
        y = Inches(3.25)
        for line in lines:
            _text(
                slide,
                line,
                left=left + Inches(0.35),
                top=y,
                width=column_w - Inches(0.7),
                height=Inches(0.62),
                size=14,
                colour=INK,
                spacing=1.25,
            )
            y += Inches(0.68)


def closing(deck: Presentation, heading: str, lines: list[str], footer: str) -> None:
    slide = blank(deck, colour=INK)

    _rule(slide, top=Inches(1.5), width=Inches(1.6))
    _text(
        slide,
        heading,
        left=MARGIN,
        top=Inches(1.9),
        width=Inches(11),
        height=Inches(1.0),
        size=42,
        bold=True,
        colour=BONE,
        font=SERIF,
    )
    y = Inches(3.2)
    for line in lines:
        _text(
            slide,
            line,
            left=MARGIN,
            top=y,
            width=Inches(10.5),
            height=Inches(0.6),
            size=19,
            colour=RGBColor(0xC9, 0xBE, 0xAC),
            spacing=1.3,
        )
        y += Inches(0.72)

    _text(
        slide,
        footer,
        left=MARGIN,
        top=Inches(6.4),
        width=Inches(11),
        height=Inches(0.5),
        size=13,
        colour=RGBColor(0x8E, 0x82, 0x74),
    )


# --------------------------------------------------------------------------
# The deck
# --------------------------------------------------------------------------
def build() -> Presentation:
    deck = Presentation()
    deck.slide_width = WIDTH
    deck.slide_height = HEIGHT

    title_slide(
        deck,
        "Stratum",
        "One platform for excavation, collections and everything\naround them — from the trench to the publication.",
        "Archaeological research and heritage management",
    )

    statement(
        deck,
        "the problem",
        "The record of an excavation is scattered\nacross nine places, and one of them\nis somebody's laptop.",
        "Context sheets in a ring binder. Finds in a spreadsheet. Photographs on an SD card, "
        "then a hard drive. The Harris matrix on graph paper, photographed. Permits in an "
        "email thread. Costs in a second spreadsheet. Nothing points at anything else.",
    )

    bullets(
        deck,
        "What that costs, in practice",
        [
            (
                "The excavation cannot be searched",
                "Answering “which contexts produced Nabataean sherds?” means opening files, "
                "not asking a question.",
            ),
            (
                "The link between object and context is fragile",
                "It lives in a column somebody typed. When the numbering changes mid-season, "
                "it quietly stops being true.",
            ),
            (
                "Only one person knows where anything is",
                "That person leaves, or is in the field with no signal, and the season stops.",
            ),
            (
                "Publication starts by rebuilding the dataset",
                "Two months of reconciling spreadsheets before a word is written — every time.",
            ),
        ],
        kicker="the problem",
    )

    statement(
        deck,
        "the idea",
        "One database. One login.\nSix ways in.",
        "Everything an institution records about a site, an object, a season or a store — in a "
        "single system that knows how the pieces relate, and that lets each person see exactly "
        "the part of it that is theirs.",
    )

    grid(
        deck,
        "Six modules, one platform",
        [
            (
                "Archaeology",
                "Projects, sites, contexts, finds. Stratigraphy, floor plans, "
                "coordinates, photographs, 3D models.",
            ),
            (
                "Museum collection",
                "Accessioning under the institution's own numbering, conservation "
                "history, exhibitions, loans, environment.",
            ),
            (
                "Office & storage",
                "Where everything physically is — from a shelf in the store to a "
                "total station signed out for the season.",
            ),
            (
                "Activity hub",
                "What was actually done: equipment, permits, preparations, costs. "
                "Repeat last year's season in one click.",
            ),
            (
                "Management",
                "Budgets, tasks, the shared calendar. Assign work; it appears on "
                "that person's dashboard.",
            ),
            (
                "Outreach",
                "Posts and channels, drafted against real records so what goes "
                "out matches what is in the archive.",
            ),
        ],
        kicker="what it is",
    )

    bullets(
        deck,
        "The part that makes it usable by an institution",
        [
            (
                "Permission is per module, not per person",
                "A collections manager needs no excavation access. A field director needs no "
                "access to the store's valuations. Neither has to be trusted with the other.",
            ),
            (
                "A restricted site stays restricted everywhere",
                "Blurred on the map, blank in the export, absent from search. Protecting a "
                "location in one place and printing it in another protects nothing.",
            ),
            (
                "Every change is attributed and reversible",
                "Who changed what, when, and what it was before. Deleting something writes a "
                "local copy first.",
            ),
            (
                "It refuses impossible data",
                "A stratigraphic sequence that loops is a typing mistake, not a discovery. It "
                "is caught on import, before it becomes a published phasing.",
            ),
        ],
        kicker="why it holds up",
    )

    compare(
        deck,
        "A season, before and after",
        [
            "Finds numbered in a spreadsheet, re-keyed into a second one for the museum.",
            "The matrix drawn on paper and redrawn in Illustrator.",
            "Photographs named by camera, matched to finds by memory.",
            "“Send me everything on Tell el-Demo” → a folder of seven CSVs.",
            "Permit renewals remembered, or not.",
        ],
        [
            "Numbered once. The museum record points at the excavation record.",
            "The matrix imported from the sheet already kept — and checked for loops.",
            "Photographs attached to the context, with their EXIF read on upload.",
            "One workbook, one sheet per kind of record, identifiers resolved to names.",
            "Permits, costs and preparations recorded against the activity, with what is outstanding.",
        ],
    )

    bullets(
        deck,
        "Built to meet data where it already is",
        [
            (
                "Import from the spreadsheets you have",
                "Column names are matched to fields — “Acc. No.”, “Inv. no.”, “Reg No” all find "
                "the same place. Every guess is shown for approval, and nothing is written until "
                "you say so.",
            ),
            (
                "A Harris matrix from a sheet",
                "Two columns and a relationship word. “Overlies”, “truncated by”, “earlier than” "
                "are all understood. A sequence that cannot have happened is refused with the "
                "loop named.",
            ),
            (
                "Legacy numbering survives",
                "A collection that cannot record 1974.1a-bis is a collection that stays in its "
                "spreadsheet. Numbers that predate the current scheme are kept and flagged, "
                "not rejected.",
            ),
            (
                "Out again as a real workbook",
                "Everything on a site or a collection in one file, with a cover sheet saying "
                "what it is, who exported it and when.",
            ),
        ],
        kicker="migration",
    )

    grid(
        deck,
        "In the field and in the store",
        [
            ("QR labels", "Printed to the sticker size you buy. A scan opens the record — "
                          "and reveals nothing the scanner could not already see."),
            ("Storage as a tree", "Building, room, cabinet, shelf, box. Move a box and "
                                  "everything in it moves with it."),
            ("Spreadsheet view", "The catalogue as a grid you tab through and paste into, "
                                 "with the platform's rules still applied to every cell."),
            ("Floor plans", "Site photographs tied to the grid they were taken in, on a "
                            "plan you can draw on."),
            ("Shared calendar", "Everybody can add to it. An event can be built from a "
                                "past activity, and brings its checklist with it."),
            ("Works offline-ish", "Runs on your own hardware, on your own network. No "
                                  "account with anybody, no subscription, no data leaving."),
        ],
        kicker="day to day",
    )

    statement(
        deck,
        "where it runs",
        "On your machine. On your network.\nOn a disk you choose.",
        "One folder and one double-click to start. Files and nightly backups can be pointed at "
        "any drive you like. Nothing is sent anywhere, no account is created with anybody, and "
        "the whole thing can be handed to another institution as a folder.",
    )

    bullets(
        deck,
        "Where it stands",
        [
            (
                "Working now",
                "All six modules, the permission model, import and export, the Harris matrix, "
                "storage, labels, the activity hub, the calendar, user administration.",
            ),
            (
                "In progress",
                "The recording sheets, pottery analysis and publication workflow as archaeologists "
                "actually use them; a reference library that links citations to sites and objects.",
            ),
            (
                "Next",
                "The digital archive — long-term preservation with checksums and a deposit "
                "record — and a request-and-upload flow for material held by other people.",
            ),
        ],
        kicker="status",
    )

    closing(
        deck,
        "What we are asking for",
        [
            "A season's worth of real records to test it against.",
            "Two people who will use it daily and say what is wrong.",
            "A decision on where the institution's copy will live.",
        ],
        "Stratum — archaeological research and heritage management",
    )

    return deck


if __name__ == "__main__":
    out = Path(__file__).resolve().parent / "Stratum.pptx"
    build().save(out)
    print(f"Wrote {out}")
